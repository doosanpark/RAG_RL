"""해결책 A 발표용 PPT 생성 (16:9, 한글).

핸드아웃 6쪽 요구사항을 슬라이드로 매핑.
첫 슬라이드는 학번/이름/github 링크 (--student / --github 인자).

실행:
    .\.venv\Scripts\python.exe -u -X utf8 -m src.sol_a.build_pptx --student "20211234 홍길동" --github "https://github.com/erid3232/RAG_RL"
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import List, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
RES = ROOT / "results"
FIG = RES / "sol_a_learning_curves.png"
CASES_MD = RES / "sol_a_qualitative_cases.md"
KFONT = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x2A, 0x4D)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)


def _set_east_asian(run, font_name: str = KFONT) -> None:
    rPr = run._r.get_or_add_rPr()
    for child in rPr.findall(qn("a:ea")):
        rPr.remove(child)
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font_name})
    rPr.append(ea)


def style_run(run, size: int = 18, bold: bool = False, color: RGBColor = None) -> None:
    run.font.name = KFONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    _set_east_asian(run, KFONT)


def add_title_only(prs: Presentation, title: str, subtitle: str = "") -> object:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # title bar
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    style_run(r, size=28, bold=True, color=NAVY)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5))
        ptf = sb.text_frame
        pr = ptf.paragraphs[0].add_run(); pr.text = subtitle
        style_run(pr, size=14, color=GRAY)
    # underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(12.3), Emu(40000))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    return slide


def add_bullets(slide, bullets: Sequence, top: float = 1.8, left: float = 0.6,
                width: float = 12.0, height: float = 5.2, size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(4)
        r = p.add_run(); r.text = ("• " if level == 0 else "– ") + str(text)
        style_run(r, size=size - 2 * level)


def add_table(slide, rows: List[List[str]], top: float = 1.8, left: float = 0.5,
              width: float = 12.3, height: float = 4.0, header_bold: bool = True) -> None:
    n_rows, n_cols = len(rows), len(rows[0])
    table = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height)).table
    for j, h in enumerate(rows[0]):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame; tf.paragraphs[0].text = ""
        r = tf.paragraphs[0].add_run(); r.text = h
        style_run(r, size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    for i, row in enumerate(rows[1:], 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            tf = cell.text_frame; tf.paragraphs[0].text = ""
            r = tf.paragraphs[0].add_run(); r.text = str(val)
            is_bold = header_bold and j == 0
            style_run(r, size=12, bold=is_bold)


def add_image(slide, image_path: Path, top: float = 1.9, left: float = 1.5, width: float = 10.0) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_code(slide, code: str, top: float = 1.8, left: float = 0.6, width: float = 12.0,
             height: float = 5.0, size: int = 12) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF7)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = "Consolas"; r.font.size = Pt(size)


def build(student: str, github: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: 표지
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "해결책 A — 검색·추론을 LLM 안에 두기"
    style_run(r, size=36, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "Qwen LoRA를 SFT→GRPO RL로 파인튜닝 (Search-R1 방식)"
    style_run(r2, size=20, color=GRAY)
    p3 = tf.add_paragraph(); p3.space_before = Pt(24)
    r3 = p3.add_run(); r3.text = "RL 클래스 프로젝트 (선행: step-wise selection RL의 후속)"
    style_run(r3, size=14, color=GRAY)
    info = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
    itf = info.text_frame
    for line, sz in [(student, 18), (github, 14)]:
        p = itf.add_paragraph() if itf.paragraphs[0].text else itf.paragraphs[0]
        rr = p.add_run(); rr.text = line
        style_run(rr, size=sz, bold=(sz == 18))

    # ---- Slide 2: 프로젝트 목표
    s = add_title_only(prs, "프로젝트 목표", "Multi-hop QA에서 RL이 무엇을 줄 수 있는가?")
    add_bullets(s, [
        "기반 프로젝트(폴백): frozen LLM + 작은 selector를 REINFORCE로 학습",
        ("→ 결과: cosine 휴리스틱(0.370) 못 넘음(0.355), transfer -24% 실패", 1),
        ("→ 진단: 추론 그릇이 없어서 multi-hop이 안 됨", 1),
        "본 프로젝트(해결책 A): 추론을 LLM 안에 넣자",
        ("Qwen2.5-0.5B를 LoRA로 SFT warmup → GRPO RL", 1),
        ("Search-R1 프로토콜로 검색·추론을 한 turn 안에서 수행", 1),
        "핵심 질문",
        ("HA-1: 추론을 LLM 안에 두면 cosine·selection-RL을 넘는가?", 1),
        ("HA-2: outcome F1 보상의 RL이 SFT 위에 추가 이득을 주는가?", 1),
        ("HA-3: HotpotQA 학습 정책이 새 도메인에 전이되는가?", 1),
    ])

    # ---- Slide 3: 환경 및 데이터셋
    s = add_title_only(prs, "환경 및 데이터셋")
    add_bullets(s, [
        "HotpotQA distractor split (in-domain)",
        ("train 90,447 / val 7,405. 질문당 후보 단락 10개(gold 2 + distractor 8)", 1),
        ("supporting_facts: dict(title, sent_id) → SFT trace 자동 생성에 사용", 1),
        ("타입: bridge(체인 추론, ~81%) / comparison(비교, ~19%)", 1),
        "Sports 룰북 QA (transfer, 350문항)",
        ("배구 등 공식 룰북을 HotpotQA 포맷으로 자체 라벨링. 도메인 전이 측정용", 1),
        "검색기: sentence-transformers/all-MiniLM-L6-v2 (384d)",
        ("set_pool로 후보 임베딩 1회 → query top-k cosine retrieve, 100단어 컷", 1),
        "Preprocessing",
        ("SFT trace 생성: gold title을 질문 단어겹침 순으로 정렬 → 멀티홉 trace 자동 작성", 1),
        ("학습/검증 분리: SFT 4000+400, RL train 4000, dev val[0:64], held-out val[200:400]", 1),
    ])

    # ---- Slide 4: 프로토콜
    s = add_title_only(prs, "Search-R1 스타일 프로토콜")
    add_code(s, (
        "<think>추론</think>\n"
        "<search>질의1</search>\n"
        "  → env가 </search>에서 generate 정지하고 다음을 주입:\n"
        "  <information>(1) Title: ...(상위 2개 passage)... </information>\n"
        "<think>추론(받은 정보 반영)</think>\n"
        "<search>질의2</search>\n"
        "  <information>...</information>\n"
        "<think>최종 결합</think>\n"
        "<answer>최종 답</answer>"
    ), size=13)
    add_bullets(s, [
        "한 assistant 턴 안에서 도구 결과가 inline 주입됨 → 추론·검색·답을 하나의 시퀀스로 학습",
        "외부 인덱스 없음. 후보 풀(10개) 내 query-based retrieve로 한정 → 무거운 엔지니어링 제거",
    ], top=5.6, height=1.5)

    # ---- Slide 5: State/Action/Reward
    s = add_title_only(prs, "State / Action / Reward 설계")
    add_bullets(s, [
        "State: (질문, 지금까지의 텍스트, 후보 풀, 현재 step)",
        ("토큰열 = 프롬프트 + (모델 생성 + 주입 info)×k + 모델 생성", 1),
        ("action_mask: 모델이 생성한 토큰만 1 → 정책경사 대상", 1),
        "Action: 매 토큰의 vocab 분포에서 샘플 (LoRA-파인튜닝된 정책)",
        ("특수 토큰 시퀀스 </search> 생성 → env가 검색 후 결과 주입 → 다음 action 계속", 1),
        ("</answer> 또는 max_turns 도달 시 종료", 1),
        "Reward (outcome rule reward, sparse):",
        ("R = answer_F1(pred, gold) + 0.1 · format_bonus", 1),
        ("format_bonus: 검색 ≥ 1회 AND <answer> 잘 형성된 경우만 +0.1", 1),
        ("F1 ∈ [0,1]이 주신호. format을 작게 둬 reward hacking 방어 (과거 교훈)", 1),
    ])

    # ---- Slide 6: 알고리즘 — 강의 범위 위치
    s = add_title_only(prs, "알고리즘: GRPO = REINFORCE + group baseline + KL",
                       subtitle="강의 06(REINFORCE+baseline) 자연 확장으로 위치")
    add_bullets(s, [
        "GRPO 한 step",
        ("1) 질문 q마다 G=5 rollout 샘플링 → 보상 r_1..r_G", 1),
        ("2) Group-normalized advantage: A_i = (r_i − mean_g) / (std_g + ε)", 1),
        ("3) 정책경사: L_PG = −E[A_i · log π(action_token)] (action 토큰만)", 1),
        ("4) KL 정규화: L = L_PG + β · KL(π || π_ref)  (π_ref = SFT 초기정책)", 1),
        ("5) Adam + grad clip 1.0으로 업데이트", 1),
        "강의에서 배운 두 요소가 핵심",
        ("Baseline(variance reduction): 강의 06의 baseline이 group-mean으로 일반화", 1),
        ("KL 정규화: 정책이 SFT 초기에서 너무 멀어지지 않게(포맷 붕괴 방어)", 1),
        "value network 불필요 → 8GB GPU 친화 (PPO 대비 메모리↓)",
    ])

    # ---- Slide 7: 학습 설정 / hyperparameters
    s = add_title_only(prs, "학습 설정 (hyperparameters)")
    add_table(s, [
        ["단계", "key 설정", "비고"],
        ["SFT warmup", "LoRA r=16 (q/k/v/o), bf16, lr=1e-4, batch=2×gradAccum 8, 3 epoch", "val_loss 0.113 수렴, info 토큰 마스킹"],
        ["GRPO 보수적", "lr=1e-5, KL coef=0.05, G=5, batch_q=4, 60 step", "policy 거의 안 움직임(dev 0.418≈SFT)"],
        ["GRPO 공격적", "lr=3e-5, KL coef=0.01, G=5, batch_q=4, 100 step", "본 실험. 3 seed(42,123,7)"],
        ["하드웨어", "RTX 4060 Ti 8GB (vLLM 없음 — Windows 미지원)", "rollout ~7s/개, ~3분/step"],
    ], top=1.7, height=2.6)
    add_bullets(s, [
        "환경 제약 대응: gradient checkpointing + max_len 768 + batch 2(logits[B,T,15만vocab] OOM 회피)",
        "dev-best 체크포인트: val[0:64] greedy F1로 매 20step 평가, peak 보존 (RL 불안정 대비)",
        "이어하기(resume): adapter+optimizer+RNG 저장 → --resume으로 끊김 복원",
    ], top=4.5, height=2.5, size=14)

    # ---- Slide 8: 실험 셋업
    s = add_title_only(prs, "실험 셋업 (평가 분리)")
    add_bullets(s, [
        "Baseline (앵커)",
        ("frozen-base: 파인튜닝 없이 같은 프로토콜로 검증 (cold-start)", 1),
        ("기존 frozen+cosine top-k: 폴백 프로젝트 보고 수치", 1),
        "Evaluation set 분리 (dev-test 누수 방지)",
        ("dev (체크포인트 선택): val[0:64], greedy rollout F1", 1),
        ("held-out (최종 보고): val[200:400] n=200, dev와 disjoint", 1),
        ("sports transfer: 350문항 전체", 1),
        "Metrics: SQuAD-style token F1 (주신호), EM, hop 타입별 (bridge/comparison)",
        "RL 3 seed (42, 123, 7) → mean ± std로 신뢰구간",
        "정직성 원칙: 단일 seed 불신, dev-best가 peak 보존, 음성 결과 정면 보고",
    ])

    # ---- Slide 9: 메인 결과 표
    s = add_title_only(prs, "실험 결과 — held-out (val[200:400] n=200, 3 seed)")
    add_table(s, [
        ["모델", "in-domain F1", "EM", "bridge", "comparison", "sports F1"],
        ["frozen-base (cold-start)", "0.006", "0.000", "0.005", "0.007", "0.005"],
        ["기존 frozen+cosine", "0.370", "—", "—", "—", "0.386"],
        ["SFT search", "0.434", "0.340", "0.435", "0.428", "0.299"],
        ["RL (3-seed mean±std)", "0.469 ± 0.007", "0.372 ± 0.002", "0.445 ± 0.010", "0.568 ± 0.024", "0.313 ± 0.023"],
    ], top=1.7, height=2.6)
    add_bullets(s, [
        "in-domain: RL > SFT > cosine. RL 이득 +0.035, seed std 0.007로 견고 → 단일 seed 노이즈 아님",
        "RL 이득은 거의 전부 comparison(+0.14). bridge는 미미(+0.01)",
        "transfer: Solution A(0.313) < 기존 cosine(0.386) — 파인튜닝의 도메인 과적합",
        "cold-start 0.006 → SFT warmup이 포맷·검색·추론 절차를 실제로 가르침을 정량 입증",
    ], top=4.7, height=2.5, size=14)

    # ---- Slide 10: 학습곡선
    s = add_title_only(prs, "학습 동역학 (dev F1, 3 seed)")
    if FIG.exists():
        add_image(s, FIG, top=1.7, left=2.6, width=8.0)
    add_bullets(s, [
        "전형적 RL 패턴: improve → peak → drift 하락",
        "seed7 step100에 dev 0.354·search 1.45로 포맷 붕괴 → dev-best가 step80 peak(0.491) 보존",
        "보수적 RL(lr 1e-5, KL 0.05)은 평평했음(여기 안 그림); 공격적(lr 3e-5, KL 0.01)에서 비로소 이득",
    ], top=5.6, height=1.7, size=13)

    # ---- Slide 11-12: 정성 케이스 (cases_md가 있으면)
    if CASES_MD.exists():
        cases_text = CASES_MD.read_text(encoding="utf-8")
        # 간단 파싱: ### 단위로 나누고 첫 2-3개만
        sections = cases_text.split("\n### ")
        case_blocks = ["### " + s_ for s_ in sections[1:4]]  # 본문 첫 3개
        # 슬라이드 1: 통계 요약 + Case 1.1
        s = add_title_only(prs, "정성 사례 (1) — RL이 SFT를 고친 비교 추론")
        add_bullets(s, [(case_blocks[0] if case_blocks else "케이스 로딩 실패", 0)],
                    top=1.8, height=5.2, size=12)
        if len(case_blocks) > 1:
            s = add_title_only(prs, "정성 사례 (2) — 둘 다 실패한 bridge (천장의 정체)")
            add_bullets(s, [(case_blocks[1], 0)], top=1.8, height=5.2, size=12)

    # ---- Slide 13: 분석
    s = add_title_only(prs, "분석 — 무엇을 의미하는가")
    add_bullets(s, [
        "HA-1 강하게 지지: 추론을 LLM 안에 두자 in-domain이 명확히 올라감 (0.370→0.434→0.469)",
        ("cold-start 0.006은 'SFT가 실제로 무엇을 가르쳤는가' 정량 증거", 1),
        "HA-2 부분 지지: RL +0.035 일관(std 0.007)이지만 이득이 comparison에 집중",
        ("comparison: 정답이 선명(yes/no, 둘 중 택1) → F1 보상이 강한 학습 신호", 1),
        ("bridge: 개방형 답 → F1 신호 약함 + SFT가 이미 retrieve를 잘함 → 헤드룸 작음", 1),
        ("⇒ RL이 고친 것은 주로 '비교 추론 결합'", 1),
        "HA-3 음성/혼합: 파인튜닝된 search 정책이 OOD에 취약",
        ("sports 0.313 < cosine 0.386 — '학습효과 vs OOD 강건성'의 트레이드오프", 1),
        ("RL은 transfer에서도 SFT보다 나음(+0.013), comparison 전이는 강함(0.507)", 1),
    ])

    # ---- Slide 14: 한계 / 향후
    s = add_title_only(prs, "한계와 향후")
    add_bullets(s, [
        "0.5B 천장: bridge가 안 오르는 핵심 이유. 1.5B 이상으로 천장↑ 필요",
        "전이 과적합: 단일 도메인 SFT의 부작용. 도메인 혼합 SFT 또는 open-retrieval로 강건성 재검토",
        "RL 비용: vLLM 부재로 rollout ~7s/개 → batched rollout 엔진 도입 시 3-5배 가속 가능",
        "SFT trace 품질: 템플릿 기반. 교사 LLM trace나 rejection sampling으로 bridge 헤드룸 확장",
        "보상 설계: F1 + format에 추가 신호(검색 품질, 근거 일치) 결합 실험",
    ])

    # ---- Slide 15: 결론
    s = add_title_only(prs, "결론")
    add_bullets(s, [
        "선행(폴백)의 'RL이 cosine을 못 넘음' 진단 — '추론 그릇 부재' — 가 정량적으로 검증됨",
        "그릇(LLM 자체 파인튜닝)을 제공하니 in-domain에서 selection-RL과 cosine을 모두 넘음",
        ("SFT가 일의 65%를 함(+0.064), RL은 그 위에 +0.035 — 주로 비교 추론을 폴리시", 1),
        ("3 seed std 0.007로 견고. 단일 seed 착시 아님(과거 교훈 통과)", 1),
        "단, 도메인 전이는 휴리스틱보다 약화 — '학습효과 vs 강건성' 트레이드오프가 정량으로 드러남",
        "정직성: 음성/혼합 결과(transfer, bridge 미개선)를 메인 표에 그대로 보고",
        ("'논문은 cherry-pick'한다는 핸드아웃 경고에 정면 대응", 1),
    ])

    out = ROOT / "report_solution_a.pptx"
    prs.save(str(out))
    print(f"[saved] {out}  ({len(prs.slides)} slides)")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--student", required=True, help="학번 이름 (예: '20211234 홍길동')")
    ap.add_argument("--github", required=True, help="github repo URL")
    args = ap.parse_args()
    build(args.student, args.github)


if __name__ == "__main__":
    main()
